"""Generate the submission deck: Ross_Agentic_Procurement.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BLUE = RGBColor(0x2A, 0x78, 0xD6)
DARK = RGBColor(0x0B, 0x0B, 0x0B)
GRAY = RGBColor(0x52, 0x51, 0x4E)
GREEN = RGBColor(0x0C, 0xA3, 0x0C)
AMBER = RGBColor(0xED, 0xA1, 0x00)
RED = RGBColor(0xD0, 0x3B, 0x3B)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def textbox(slide, left, top, width, height, text, size=18, bold=False,
            color=DARK, align=PP_ALIGN.LEFT, bullets=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    items = bullets if bullets is not None else [text]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " + item) if bullets is not None else item
        p.alignment = align
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return box


def header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    textbox(slide, 0.5, 0.18, 12.3, 0.6, title, size=28, bold=True,
            color=RGBColor(255, 255, 255))
    if subtitle:
        textbox(slide, 0.5, 0.72, 12.3, 0.4, subtitle, size=14,
                color=RGBColor(0xDC, 0xE8, 0xF8))


def chip(slide, left, top, width, height, title, body, fill=LIGHT, title_color=BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = RGBColor(0xD5, 0xDB, 0xE4); sh.line.width = Pt(0.75)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.text = title
    for r in p.runs:
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = title_color
        r.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.text = body
    for r in p2.runs:
        r.font.size = Pt(12); r.font.color.rgb = GRAY; r.font.name = "Calibri"


# ------------------------------------------------------------- 1. title
s = add_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), prs.slide_width, Inches(2.6))
band.fill.solid(); band.fill.fore_color.rgb = BLUE; band.line.fill.background()
textbox(s, 0.8, 2.55, 11.7, 0.9, "Agentic AI Procurement Advisor for Ross Stores",
        size=40, bold=True, color=RGBColor(255, 255, 255))
textbox(s, 0.8, 3.65, 11.7, 0.6,
        "When, how much, and whether to buy — under tariffs, trade deals and demand trends",
        size=20, color=RGBColor(0xDC, 0xE8, 0xF8))
textbox(s, 0.8, 5.4, 11.7, 0.5,
        "Solo project  ·  Multi-agent pipeline on Google Gemini  ·  Streamlit prototype",
        size=16, color=GRAY)

# --------------------------------------------- 2. company & value chain
s = add_slide()
header(s, "Company & Value-Chain Application",
       "Problem Statement Identification — where agentic AI fits at Ross")
textbox(s, 0.5, 1.45, 12.3, 0.5, "Ross Stores, Inc. — 'Ross Dress for Less'",
        size=22, bold=True, color=BLUE)
textbox(s, 0.5, 2.0, 6.2, 4.5, "", bullets=[
    "US off-price retailer: ~2,200 stores, ~$21B revenue",
    "Business model = opportunistic buying: procure branded goods below wholesale, resell 20–60% under department-store prices",
    "Margins live or die on procurement decisions",
    "Value-chain function chosen: Sourcing & Procurement (Operations / Buying)",
], size=16)
chip(s, 7.0, 2.0, 5.8, 1.5, "Why this function?",
     "Hundreds of buyers make thousands of buy/no-buy calls per season across brands and origin countries — the most decision-dense function in the company.")
chip(s, 7.0, 3.7, 5.8, 1.5, "What changed in 2025–26?",
     "Tariff escalations, renegotiated trade deals, and fast-moving demand trends mean yesterday's profitable buy is today's loss-maker.")

# --------------------------------------------------- 3. problem statement
s = add_slide()
header(s, "Problem Statement")
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5),
                         Inches(12.1), Inches(1.9))
box.fill.solid(); box.fill.fore_color.rgb = LIGHT
box.line.color.rgb = BLUE; box.line.width = Pt(1.5)
tf = box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.18)
p = tf.paragraphs[0]
p.text = ("“When, and in what quantity, should Ross procure a product from a given "
          "brand and origin country so the buy stays profitably feasible — given current "
          "tariffs, geopolitical trade deals and demand trends? And if the answer is "
          "‘partially yes’, exactly how much?”")
for r in p.runs:
    r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = DARK; r.font.name = "Calibri"
textbox(s, 0.6, 3.7, 12.1, 0.4, "Every candidate buy must answer four questions:",
        size=16, bold=True, color=GRAY)
chip(s, 0.6, 4.2, 2.9, 2.3, "1 · Feasibility",
     "Does landed cost (FOB + duty + freight + handling) still clear the target gross margin at the off-price retail point?")
chip(s, 3.7, 4.2, 2.9, 2.3, "2 · Risk",
     "Is the origin's US trade relationship stable, deteriorating (tariff hikes) or improving (relief being ratified)?")
chip(s, 6.8, 4.2, 2.9, 2.3, "3 · Demand",
     "Is the category trending up or down with consumers? What does next quarter's forecast support?")
chip(s, 9.9, 4.2, 2.9, 2.3, "4 · Quantity & timing",
     "Full buy, PARTIAL buy (cap exposure at MOQ+), or walk away — and order now or after a known event?")

# ------------------------------------------------------ 4. why agentic
s = add_slide()
header(s, "Why Agentic AI — not a dashboard, not one prompt")
rows = [
    ("Four heterogeneous evidence streams",
     "Tariff schedules, trade-deal news, demand indices and unit economics each need different data and reasoning → one specialist agent per stream, like a real buying office."),
    ("Synthesis with graded judgment",
     "The output isn't a score — it's 'PARTIAL YES, 8,000 units, wait for the October ratification vote'. A Decision Agent argues from four reports; a rule engine can't handle novel geopolitics."),
    ("Conditions change weekly",
     "New tariff tranches, port congestion, deal votes. Agents re-run on fresh data with zero re-engineering."),
    ("Auditability",
     "Hard numbers (landed cost, margins, forecasts) are deterministic Python; the LLM adds interpretation. Every agent's report is shown to the buyer."),
]
y = 1.5
for title, body in rows:
    chip(s, 0.6, y, 12.1, 1.32, title, body)
    y += 1.45

# ------------------------------------------------------- 5. data collation
s = add_slide()
header(s, "Data Collation — 6 synthetic datasets (seeded & reproducible)",
       "Generated by data_gen.py, calibrated to a mid-2026 trade environment")
d = [
    ("brand_catalog.csv", "24 candidate brands · 8 origin countries · 5 categories · FOB cost, MOQ, lead time, quality, freight"),
    ("tariff_schedule.csv", "Base vs current US tariff per country × category — reflects 2025–26 tariff actions (China +30pts, Mexico duty-free under USMCA)"),
    ("trade_agreements.csv", "Per-country deal status: active / suspended / under negotiation, relief terms, outlook"),
    ("geopolitical_events.csv", "News-style event feed with 1–10 risk impact (tariff tranches, port congestion, ratification votes, currency moves)"),
    ("market_trends.csv", "24 months × 5 categories: demand index, social buzz, retail price points"),
    ("sales_history.csv", "24 months of Ross category sales (units, revenue) with seasonality"),
]
y = 1.5
for i, (name, body) in enumerate(d):
    chip(s, 0.6 + (i % 2) * 6.25, y + (i // 2) * 1.85, 5.9, 1.7, name, body)

# ------------------------------------------------- 6. solution architecture
s = add_slide()
header(s, "Solution Architecture — the agent team", "agents.py · Google Gemini · JSON-constrained outputs")
agents = [
    ("Tariff Agent", "Duty exposure, landed-cost breakdown, cheaper-origin alternatives", 0.5),
    ("Geopolitics Agent", "Trade-deal status, country risk 1–10, 6-month outlook", 3.7),
    ("Trend Agent", "Demand momentum, seasonality, next-quarter forecast", 6.9),
    ("Finance Agent", "Margin vs target, breakeven retail, negotiation levers", 10.1),
]
for name, body, x in agents:
    chip(s, x, 1.6, 2.9, 1.9, name, body)
arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(3.6), Inches(0.9), Inches(0.8))
arrow.fill.solid(); arrow.fill.fore_color.rgb = BLUE; arrow.line.fill.background()
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(4.5), Inches(7.0), Inches(1.6))
sh.fill.solid(); sh.fill.fore_color.rgb = BLUE; sh.line.fill.background()
tf = sh.text_frame; tf.word_wrap = True; tf.margin_top = Inches(0.15)
p = tf.paragraphs[0]; p.text = "Chief Decision Agent"; p.alignment = PP_ALIGN.CENTER
for r in p.runs:
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = RGBColor(255, 255, 255); r.font.name = "Calibri"
p2 = tf.add_paragraph()
p2.text = "YES / PARTIAL_YES / NO · recommended quantity · order timing · confidence · revisit triggers"
p2.alignment = PP_ALIGN.CENTER
for r in p2.runs:
    r.font.size = Pt(13); r.font.color.rgb = RGBColor(0xDC, 0xE8, 0xF8); r.font.name = "Calibri"
textbox(s, 3.2, 6.3, 7.0, 0.6,
        "Guardrails: quantity clamped to [MOQ, requested] in multiples of 500; forced to 0 on NO",
        size=13, color=GRAY, align=PP_ALIGN.CENTER)

# --------------------------------------------------------- 7. how it decides
s = add_slide()
header(s, "How a decision is made", "Deterministic math + LLM judgment = auditable recommendations")
steps = [
    ("1 · Buyer input", "Category, brand, requested quantity, target margin — from the Streamlit sidebar"),
    ("2 · Deterministic core", "Landed cost = FOB + duty + freight + handling; gross margin vs retail; next-quarter forecast from sales drift"),
    ("3 · Specialist agents", "Each gets only its data slice + the hard numbers; returns strict JSON (risk levels, outlooks, viability)"),
    ("4 · Decision synthesis", "Chief agent weighs all four reports → decision, quantity, timing, rationale, conditions"),
    ("5 · Buyer-facing output", "Decision card, cost/tariff/trend charts, and each agent's full report for audit"),
]
y = 1.5
for title, body in steps:
    chip(s, 0.6, y, 12.1, 1.02, title, body)
    y += 1.13

# ------------------------------------------------------- 8. demo scenarios
s = add_slide()
header(s, "Prototype in action — three demo scenarios")
chip(s, 0.6, 1.6, 3.95, 4.3, "China · MeiTex Garments (Apparel)",
     "Tariff 16.5% → 46.5% after 2026 tranches. Duty alone adds ~$2.37/unit on a $5.10 FOB. "
     "Margin collapses below target → typically NO or a sharply reduced buy, with 'revisit if talks resume'.",
     title_color=RED)
chip(s, 4.7, 1.6, 3.95, 4.3, "Vietnam · DaNang Soles (Footwear)",
     "Footwear demand rising; tariff relief (−10pts) awaiting an October ratification vote. "
     "Typical outcome: PARTIAL_YES — order near MOQ now to secure capacity, place the balance after the vote.",
     title_color=AMBER)
chip(s, 8.8, 1.6, 3.95, 4.3, "Mexico · Casa Loma Living (Home Goods)",
     "USMCA duty-free, 12-day lead time, stable politics. Landed cost stays low → clean YES at or near requested quantity, order now.",
     title_color=GREEN)
textbox(s, 0.6, 6.1, 12.1, 0.6,
        "Same pipeline, three coherent, differentiated answers — driven by tariffs, deals and trends, not hand-tuned rules.",
        size=15, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# ------------------------------------------------------------ 9. tech stack
s = add_slide()
header(s, "Tech Stack & Repository")
chip(s, 0.6, 1.6, 5.9, 1.5, "LLM", "Google Gemini (gemini-2.5-flash) via google-genai SDK — JSON-mode, temperature 0.2, model fallback")
chip(s, 0.6, 3.25, 5.9, 1.5, "Orchestration", "Plain-Python sequential pipeline (agents.py) — transparent, explainable, no framework overhead")
chip(s, 0.6, 4.9, 5.9, 1.5, "UI", "Streamlit dashboard (app.py) with Plotly charts: tariff comparison, demand trends, landed-cost breakdown")
chip(s, 7.0, 1.6, 5.8, 1.5, "Data", "Six seeded synthetic datasets (data_gen.py) — reproducible, mid-2026 calibrated")
chip(s, 7.0, 3.25, 5.8, 1.5, "Guardrails", "Deterministic landed-cost & margin math; quantity clamped to MOQ/request; decision forced consistent")
chip(s, 7.0, 4.9, 5.8, 1.5, "Run it", "pip install -r requirements.txt · add GEMINI_API_KEY to .env · streamlit run app.py")

# ------------------------------------------------------- 10. conclusion
s = add_slide()
header(s, "Impact & Future Scope")
textbox(s, 0.5, 1.5, 12.3, 0.5, "What the prototype demonstrates", size=20, bold=True, color=BLUE)
textbox(s, 0.5, 2.05, 12.3, 2.0, "", bullets=[
    "Turns a multi-week merchant judgment call into a 30-second, auditable, evidence-backed recommendation",
    "Answers the 'partial yes' question with an exact unit count and timing — not just a red/green light",
    "Buyers can flip origins (China → Vietnam → Mexico) and watch the recommendation shift coherently",
], size=16)
textbox(s, 0.5, 4.1, 12.3, 0.5, "Future scope", size=20, bold=True, color=BLUE)
textbox(s, 0.5, 4.65, 12.3, 2.2, "", bullets=[
    "Live data: replace synthetic feeds with USITC tariff APIs, news APIs and Google Trends",
    "Portfolio mode: optimize across all candidate buys under a total open-to-buy budget",
    "Autonomous monitoring: agents watch the event feed and proactively alert buyers when a placed order's assumptions break",
], size=16)
textbox(s, 0.5, 6.8, 12.3, 0.5, "Thank you", size=18, bold=True, color=GRAY,
        align=PP_ALIGN.CENTER)

prs.save("Ross_Agentic_Procurement.pptx")
print(f"Saved Ross_Agentic_Procurement.pptx with {len(prs.slides.slides if hasattr(prs.slides, 'slides') else prs.slides._sldIdLst)} slides")
